"""
Accenture-branded deck: Billing Reconciliation Agent — Problem Statement,
ROI Logic, MVP Boundary, Go Decision.

Brand rules applied:
- Core purple #A100FF, purple spectrum, neutrals only
- Graphik font (Arial fallback)
- RECTANGLE shapes only — no rounded corners
- Cover: full Accenture logo
- Internal slides: Greater Than (GT) symbol top-right
- Light mode default; sentence case headings
- 60-70% neutrals, 30-40% purple, <5% secondary
"""

from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

SKILL_DIR = Path("/Users/preetinder.sidhu/RDEProjects/capstone-pod-4-northwind/.claude/skills/creating-accenture-presentations")
LOGO_PATH  = SKILL_DIR / "assets/logos/Acc_Logo_Black_Purple_RGB.png"
GT_PATH    = SKILL_DIR / "assets/greater-than/Acc_GT_Solid_P1_RGB.png"
OUT_PATH   = Path("/Users/preetinder.sidhu/RDEProjects/FDE-Meridian-Pod1/reconciliation-agent-deck.pptx")

# ── Brand palette ────────────────────────────────────────────────────────────
PURPLE      = RGBColor(0xA1, 0x00, 0xFF)   # core purple
DARK_PURPLE = RGBColor(0x75, 0x00, 0xC0)
DEEP_PURPLE = RGBColor(0x46, 0x00, 0x73)
LIGHT_PURPLE= RGBColor(0xE6, 0xDC, 0xFF)
BLACK       = RGBColor(0x00, 0x00, 0x00)
DARK_GRAY   = RGBColor(0x81, 0x81, 0x80)
MID_GRAY    = RGBColor(0xCF, 0xCF, 0xCF)
LIGHT_GRAY  = RGBColor(0xF1, 0xF1, 0xEF)
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)

FONT = "Graphik"

# Slide size: 16x9 = 10" x 5.625"
W = Inches(10)
H = Inches(5.625)


def new_prs():
    prs = Presentation()
    prs.slide_width  = W
    prs.slide_height = H
    return prs


def blank_slide(prs):
    layout = prs.slide_layouts[6]  # completely blank
    return prs.slides.add_slide(layout)


def add_text(slide, text, x, y, w, h, size=14, bold=False, color=BLACK,
             align=PP_ALIGN.LEFT, wrap=True, font=FONT):
    txb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf  = txb.text_frame
    tf.word_wrap = wrap
    p   = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name  = font
    run.font.size  = Pt(size)
    run.font.bold  = bold
    run.font.color.rgb = color


def add_rect(slide, x, y, w, h, fill=PURPLE, line_color=None, line_width=0):
    from pptx.util import Pt as Pt2
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE = 1
        Inches(x), Inches(y), Inches(w), Inches(h)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = Pt2(line_width)
    else:
        shape.line.fill.background()
    return shape


def add_gt_symbol(slide):
    if GT_PATH.exists():
        slide.shapes.add_picture(str(GT_PATH), Inches(9.3), Inches(0.15), Inches(0.55), Inches(0.55))


def add_rich_text(slide, items, x, y, w, h, default_size=12, default_color=BLACK):
    """items: list of (text, bold, size, color, is_bullet)"""
    from pptx.util import Inches as I2, Pt as P2
    from pptx.oxml.ns import qn
    import lxml.etree as etree

    txb = slide.shapes.add_textbox(I2(x), I2(y), I2(w), I2(h))
    tf  = txb.text_frame
    tf.word_wrap = True

    first = True
    for (text, bold, size, color, is_bullet) in items:
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()

        if is_bullet:
            pPr = p._pPr if p._pPr is not None else p._p.get_or_add_pPr()
            buChar = etree.SubElement(pPr, qn("a:buChar"))
            buChar.set("char", "•")
            indent = etree.SubElement(pPr, qn("a:buFont"))
            indent.set("typeface", FONT)
            p.level = 0

        run = p.add_run()
        run.text = text
        run.font.name  = FONT
        run.font.size  = Pt(size or default_size)
        run.font.bold  = bold
        run.font.color.rgb = color or default_color


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 1 — Cover
# ─────────────────────────────────────────────────────────────────────────────
def slide_cover(prs):
    sl = blank_slide(prs)

    # Full white background
    add_rect(sl, 0, 0, 10, 5.625, fill=WHITE)

    # Left purple block
    add_rect(sl, 0, 0, 0.25, 5.625, fill=PURPLE)

    # Bottom purple strip
    add_rect(sl, 0, 4.9, 10, 0.725, fill=DEEP_PURPLE)

    # Accenture logo top-left
    if LOGO_PATH.exists():
        sl.shapes.add_picture(str(LOGO_PATH), Inches(0.45), Inches(0.3), Inches(2.2), Inches(0.6))

    # Main title
    add_text(sl, "Billing reconciliation agent",
             0.45, 1.4, 8.5, 0.9, size=32, bold=True, color=BLACK)

    # Purple highlight bar behind subtitle phrase
    add_rect(sl, 0.45, 2.35, 5.1, 0.48, fill=PURPLE)
    add_text(sl, "from manual effort to intelligent exception triage",
             0.45, 2.35, 5.8, 0.48, size=14, bold=False, color=WHITE)

    # Subtitle line
    add_text(sl, "PRJ-NS-7421  |  Northstar Civic Group  |  Cycle 2026-04",
             0.45, 2.95, 8, 0.4, size=11, color=DARK_GRAY)

    # Bottom strip text
    add_text(sl, "Pod 1  |  Meridian Atlas Partners  |  Hackathon Day 2",
             0.45, 5.0, 7, 0.35, size=10, color=WHITE)


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 2 — Problem statement
# ─────────────────────────────────────────────────────────────────────────────
def slide_problem(prs):
    sl = blank_slide(prs)
    add_rect(sl, 0, 0, 10, 5.625, fill=WHITE)
    add_rect(sl, 0, 0, 0.12, 5.625, fill=PURPLE)
    add_gt_symbol(sl)

    # Heading
    add_text(sl, "Problem statement", 0.35, 0.2, 9, 0.45, size=20, bold=True, color=BLACK)
    add_rect(sl, 0.35, 0.7, 9.3, 0.04, fill=MID_GRAY)

    # One-sentence definition — purple highlight
    add_rect(sl, 0.35, 0.85, 9.3, 0.55, fill=LIGHT_PURPLE)
    add_text(sl,
        "An agentic pipeline that reads every backup document line by line, applies contract "
        "rules and institutional memory to classify and auto-resolve expense exceptions, and "
        "presents a pre-populated decision queue to the analyst.",
        0.45, 0.88, 9.1, 0.52, size=10.5, bold=False, color=DEEP_PURPLE)

    # Three problem columns
    col_data = [
        ("Manual volume",
         "~27 expense transactions per cycle. Analyst opens every receipt, checks every amount, applies every cap by hand."),
        ("Fragmented memory",
         "Prior decisions live in email, Teams, and individual recall. No structured store. Instructions contradict each other across cycles."),
        ("False CLEAN risk",
         "Document totals can match SAP while containing non-reimbursable items. A dinner receipt may include alcohol. A total-level check misses it."),
    ]

    for i, (title, body) in enumerate(col_data):
        x = 0.35 + i * 3.18
        add_rect(sl, x, 1.58, 3.0, 0.42, fill=DEEP_PURPLE)
        add_text(sl, title, x + 0.1, 1.6, 2.8, 0.38, size=11.5, bold=True, color=WHITE)
        add_text(sl, body,  x + 0.1, 2.1, 2.8, 1.1, size=10, color=BLACK)

    # Current state vs opportunity
    add_rect(sl, 0.35, 3.35, 4.4, 0.35, fill=MID_GRAY)
    add_text(sl, "Current state", 0.45, 3.37, 4.2, 0.3, size=10.5, bold=True, color=BLACK)

    current_items = [
        "Analyst reviews all transactions regardless of complexity",
        "Alcohol, FX errors, missing markup caught only if analyst remembers",
        "No audit trail — decisions live in personal notes",
    ]
    for j, item in enumerate(current_items):
        add_text(sl, f"  {item}", 0.35, 3.75 + j * 0.28, 4.4, 0.28, size=9.5, color=BLACK)

    add_rect(sl, 5.05, 3.35, 4.6, 0.35, fill=PURPLE)
    add_text(sl, "Opportunity", 5.15, 3.37, 4.4, 0.3, size=10.5, bold=True, color=WHITE)

    opp_items = [
        "Agent handles classification; analyst reviews exceptions only",
        "Contract rules + prior resolutions applied on every run",
        "Immutable audit trail, 2-year retention",
    ]
    for j, item in enumerate(opp_items):
        add_text(sl, f"  {item}", 5.05, 3.75 + j * 0.28, 4.6, 0.28, size=9.5, color=BLACK)


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 3 — ROI Logic
# ─────────────────────────────────────────────────────────────────────────────
def slide_roi(prs):
    sl = blank_slide(prs)
    add_rect(sl, 0, 0, 10, 5.625, fill=WHITE)
    add_rect(sl, 0, 0, 0.12, 5.625, fill=PURPLE)
    add_gt_symbol(sl)

    add_text(sl, "ROI logic", 0.35, 0.2, 9, 0.45, size=20, bold=True, color=BLACK)
    add_rect(sl, 0.35, 0.7, 9.3, 0.04, fill=MID_GRAY)

    # Two columns: Current vs Future
    add_rect(sl, 0.35, 0.85, 4.35, 0.38, fill=DARK_GRAY)
    add_text(sl, "Current state — every cycle", 0.45, 0.87, 4.15, 0.34, size=11, bold=True, color=WHITE)

    add_rect(sl, 5.0, 0.85, 4.65, 0.38, fill=PURPLE)
    add_text(sl, "Future state — with agent", 5.1, 0.87, 4.45, 0.34, size=11, bold=True, color=WHITE)

    current = [
        "Analyst opens all 27 expense transactions",
        "Manually checks each receipt amount",
        "Applies FX, markup, caps from memory",
        "Chases missing receipts by email",
        "No structure — prior decisions recalled by individual",
        "False CLEAN risk on every document",
    ]
    future = [
        "Agent classifies all 27 in under 5 minutes",
        "Analyst reviews flagged queue only (~13 items)",
        "Contract rules applied automatically on every run",
        "Missing docs surfaced instantly as MISSING_DOC",
        "Exception store — structured, project-scoped, auditable",
        "Line-level extraction eliminates false CLEAN",
    ]

    for j, (c, f) in enumerate(zip(current, future)):
        y = 1.35 + j * 0.38
        # current — light gray row alternating
        bg = LIGHT_GRAY if j % 2 == 0 else WHITE
        add_rect(sl, 0.35, y, 4.35, 0.35, fill=bg)
        add_text(sl, c, 0.45, y + 0.04, 4.2, 0.3, size=9.5, color=BLACK)
        # future — light purple row alternating
        bg2 = LIGHT_PURPLE if j % 2 == 0 else WHITE
        add_rect(sl, 5.0, y, 4.65, 0.35, fill=bg2)
        add_text(sl, f, 5.1, y + 0.04, 4.5, 0.3, size=9.5, color=DEEP_PURPLE if j % 2 == 0 else BLACK)

    # Arrow between columns
    add_rect(sl, 4.72, 2.3, 0.24, 0.24, fill=PURPLE)
    add_text(sl, ">", 4.73, 2.28, 0.24, 0.28, size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    # Bottom metric bar
    add_rect(sl, 0.35, 5.1, 9.3, 0.38, fill=DEEP_PURPLE)

    metrics = [
        ("55%+", "exceptions auto-resolved"),
        ("0", "false CLEANs"),
        ("100%", "decisions audited"),
        ("< 5 min", "to first exception queue"),
    ]
    for i, (val, label) in enumerate(metrics):
        x = 0.6 + i * 2.3
        add_text(sl, val,   x, 5.1,  1.5, 0.22, size=13, bold=True, color=PURPLE, align=PP_ALIGN.CENTER)
        add_text(sl, label, x, 5.3,  1.8, 0.18, size=8,  color=WHITE,  align=PP_ALIGN.CENTER)


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 4 — MVP Boundary
# ─────────────────────────────────────────────────────────────────────────────
def slide_mvp(prs):
    sl = blank_slide(prs)
    add_rect(sl, 0, 0, 10, 5.625, fill=WHITE)
    add_rect(sl, 0, 0, 0.12, 5.625, fill=PURPLE)
    add_gt_symbol(sl)

    add_text(sl, "MVP boundary", 0.35, 0.2, 9, 0.45, size=20, bold=True, color=BLACK)
    add_rect(sl, 0.35, 0.7, 9.3, 0.04, fill=MID_GRAY)

    # IN scope column
    add_rect(sl, 0.35, 0.85, 4.5, 0.38, fill=DEEP_PURPLE)
    add_text(sl, "In scope", 0.45, 0.87, 4.3, 0.34, size=12, bold=True, color=WHITE)

    in_scope = [
        ("Expense transactions", "PRJ-NS-7421, cycle 2026-04 — 27 transactions"),
        ("All document types", "Receipts, hotel folios, vendor invoices, mileage logs"),
        ("Contract rule engine", "Caps, alcohol exclusion, FX, subcontractor markup"),
        ("Prior exception matching", "7 recurring resolutions, project-scoped only"),
        ("Analyst review UI", "Streamlit — exception queue with pre-populated decisions"),
        ("Audit trail", "Immutable JSON, 2-year retention, every decision logged"),
    ]

    for j, (title, detail) in enumerate(in_scope):
        y = 1.33 + j * 0.52
        add_rect(sl, 0.35, y, 0.06, 0.38, fill=PURPLE)
        add_text(sl, title,  0.48, y,        4.2, 0.22, size=10.5, bold=True,  color=BLACK)
        add_text(sl, detail, 0.48, y + 0.22, 4.2, 0.22, size=9,   bold=False, color=DARK_GRAY)

    # OUT scope column
    add_rect(sl, 5.15, 0.85, 4.5, 0.38, fill=MID_GRAY)
    add_text(sl, "Out of scope — MVP", 5.25, 0.87, 4.3, 0.34, size=12, bold=True, color=BLACK)

    out_scope = [
        ("Labour reconciliation", "Different risk profile — rate-table lookup, not doc intelligence"),
        ("SAP write-back", "Financial data — no unauthorised changes; analyst approves first"),
        ("Teams / email ingestion", "Exception store solves the problem more reliably"),
        ("Multi-project batch", "Single project validates the pattern; Ingest agent generalises"),
        ("Production database", "JSON files sufficient for demo; production would add persistence"),
    ]

    for j, (title, detail) in enumerate(out_scope):
        y = 1.33 + j * 0.6
        add_rect(sl, 5.15, y, 0.06, 0.45, fill=DARK_GRAY)
        add_text(sl, title,  5.28, y,        4.2, 0.25, size=10.5, bold=True,  color=BLACK)
        add_text(sl, detail, 5.28, y + 0.25, 4.2, 0.22, size=9,   bold=False, color=DARK_GRAY)


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 5 — Go Decision
# ─────────────────────────────────────────────────────────────────────────────
def slide_go(prs):
    sl = blank_slide(prs)
    add_rect(sl, 0, 0, 10, 5.625, fill=WHITE)
    add_rect(sl, 0, 0, 0.12, 5.625, fill=PURPLE)
    add_gt_symbol(sl)

    add_text(sl, "Go decision", 0.35, 0.2, 9, 0.45, size=20, bold=True, color=BLACK)
    add_rect(sl, 0.35, 0.7, 9.3, 0.04, fill=MID_GRAY)

    # GO badge
    add_rect(sl, 0.35, 0.85, 1.6, 0.9, fill=PURPLE)
    add_text(sl, "GO", 0.35, 0.88, 1.6, 0.9, size=38, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    add_text(sl, "Prototype viable on synthetic dataset.\nArchitecture generalises without re-work.",
             2.1, 0.9, 7.5, 0.8, size=12, color=BLACK)

    add_rect(sl, 0.35, 1.85, 9.3, 0.04, fill=LIGHT_PURPLE)

    # Three GO reasons
    reasons = [
        (
            "55% KPI met from day one",
            "6 of 7 recurring prior resolutions map directly to exceptions in cycle 2026-04. "
            "The institutional memory already exists — the agent reads it. "
            "No new data collection required before the first run."
        ),
        (
            "The hardest problem is solved",
            "Line-level document extraction via Claude catches false CLEANs that no rule-table can catch. "
            "RC-003 hotel folio and RC-007 team dinner both contain alcohol hidden inside a matching total. "
            "Both are surfaced and flagged automatically."
        ),
        (
            "Architecture generalises without re-work",
            "Swapping the Ingest agent's data source adds a new project. "
            "The Classify and Triage agents are untouched. "
            "Labour reconciliation, multi-project batch, and SAP integration "
            "are additive — they do not require redesign."
        ),
    ]

    for i, (heading, body) in enumerate(reasons):
        x = 0.35 + i * 3.18
        # Numbered purple circle (using rectangle — no circles per brand rules)
        add_rect(sl, x, 2.05, 0.38, 0.38, fill=PURPLE)
        add_text(sl, str(i + 1), x, 2.05, 0.38, 0.38, size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_text(sl, heading, x + 0.48, 2.08, 2.6, 0.32, size=11, bold=True, color=BLACK)
        add_rect(sl, x, 2.5, 3.05, 0.04, fill=LIGHT_PURPLE)
        add_text(sl, body, x, 2.6, 3.05, 1.5, size=9.5, color=BLACK)

    # Evidence row
    add_rect(sl, 0.35, 4.35, 9.3, 0.38, fill=DEEP_PURPLE)
    add_text(sl, "Evidence", 0.45, 4.37, 1.5, 0.34, size=11, bold=True, color=WHITE)

    evidence = [
        "37 automated tests — all passing",
        "Pipeline runs on full 2026-04 dataset",
        "Streamlit UI deployed at localhost:8501",
        "Audit trail verified end-to-end",
    ]
    for i, ev in enumerate(evidence):
        x = 2.1 + i * 1.88
        add_text(sl, ev, x, 4.37, 1.82, 0.34, size=9, color=WHITE)

    # Footer
    add_rect(sl, 0.35, 4.88, 9.3, 0.04, fill=MID_GRAY)
    add_text(sl, "SAP remains the system of record. All work on synthetic data. No real client data used.",
             0.35, 4.95, 9.3, 0.3, size=8, color=DARK_GRAY)


# ─────────────────────────────────────────────────────────────────────────────
# Build
# ─────────────────────────────────────────────────────────────────────────────
def build():
    prs = new_prs()
    slide_cover(prs)
    slide_problem(prs)
    slide_roi(prs)
    slide_mvp(prs)
    slide_go(prs)
    prs.save(str(OUT_PATH))
    print(f"Saved: {OUT_PATH}")


if __name__ == "__main__":
    build()
