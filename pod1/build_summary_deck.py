"""
Accenture-branded deck: Reconciliation Agent — Solution Summary
9 slides covering all sections of SOLUTION-SUMMARY.md

Brand rules:
- Core purple #A100FF, purple spectrum, neutrals only
- Graphik font (Arial fallback)
- RECTANGLE shapes only
- Cover: full Accenture logo | Internal: GT symbol top-right
- Light mode, sentence case headings
"""

from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

SKILL_DIR = Path("/Users/preetinder.sidhu/RDEProjects/capstone-pod-4-northwind/.claude/skills/creating-accenture-presentations")
LOGO_PATH = SKILL_DIR / "assets/logos/Acc_Logo_Black_Purple_RGB.png"
GT_PATH   = SKILL_DIR / "assets/greater-than/Acc_GT_Solid_P1_RGB.png"
OUT_PATH  = Path("/Users/preetinder.sidhu/RDEProjects/FDE-Meridian-Pod1/solution-summary-deck.pptx")

# ── Brand palette ─────────────────────────────────────────────────────────────
PURPLE       = RGBColor(0xA1, 0x00, 0xFF)
DARK_PURPLE  = RGBColor(0x75, 0x00, 0xC0)
DEEP_PURPLE  = RGBColor(0x46, 0x00, 0x73)
LIGHT_PURPLE = RGBColor(0xE6, 0xDC, 0xFF)
BLACK        = RGBColor(0x00, 0x00, 0x00)
DARK_GRAY    = RGBColor(0x81, 0x81, 0x80)
MID_GRAY     = RGBColor(0xCF, 0xCF, 0xCF)
LIGHT_GRAY   = RGBColor(0xF1, 0xF1, 0xEF)
WHITE        = RGBColor(0xFF, 0xFF, 0xFF)

FONT = "Graphik"
W, H = Inches(10), Inches(5.625)


# ── Primitives ────────────────────────────────────────────────────────────────

def new_prs():
    prs = Presentation()
    prs.slide_width  = W
    prs.slide_height = H
    return prs


def blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def rect(sl, x, y, w, h, fill=WHITE, line=None, lw=0):
    sh = sl.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill
    if line:
        sh.line.color.rgb = line
        sh.line.width = Pt(lw)
    else:
        sh.line.fill.background()
    return sh


def txt(sl, text, x, y, w, h, size=11, bold=False, color=BLACK,
        align=PP_ALIGN.LEFT, italic=False, wrap=True):
    tb = sl.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = wrap
    p  = tf.paragraphs[0]
    p.alignment = align
    r  = p.add_run()
    r.text = text
    r.font.name   = FONT
    r.font.size   = Pt(size)
    r.font.bold   = bold
    r.font.italic = italic
    r.font.color.rgb = color


def heading(sl, text, y=0.2):
    txt(sl, text, 0.35, y, 9.0, 0.45, size=20, bold=True, color=BLACK)
    rect(sl, 0.35, y + 0.5, 9.3, 0.04, fill=MID_GRAY)


def gt(sl):
    if GT_PATH.exists():
        sl.shapes.add_picture(str(GT_PATH), Inches(9.3), Inches(0.15),
                               Inches(0.55), Inches(0.55))


def chrome(sl, heading_text):
    rect(sl, 0, 0, 10, 5.625, fill=WHITE)
    rect(sl, 0, 0, 0.12, 5.625, fill=PURPLE)
    gt(sl)
    heading(sl, heading_text)


def pill(sl, x, y, w, h, label, label_size=9, bg=PURPLE, fg=WHITE):
    rect(sl, x, y, w, h, fill=bg)
    txt(sl, label, x + 0.06, y + 0.03, w - 0.1, h - 0.06,
        size=label_size, bold=True, color=fg, align=PP_ALIGN.CENTER)


# ── Slide 1 — Cover ───────────────────────────────────────────────────────────

def slide_cover(prs):
    sl = blank(prs)
    rect(sl, 0, 0, 10, 5.625, fill=WHITE)
    rect(sl, 0, 0, 0.25, 5.625, fill=PURPLE)
    rect(sl, 0, 4.9, 10, 0.725, fill=DEEP_PURPLE)

    if LOGO_PATH.exists():
        sl.shapes.add_picture(str(LOGO_PATH), Inches(0.45), Inches(0.25),
                               Inches(2.2), Inches(0.6))

    txt(sl, "Reconciliation agent", 0.45, 1.3, 9, 0.85, size=34, bold=True, color=BLACK)

    rect(sl, 0.45, 2.28, 6.5, 0.5, fill=PURPLE)
    txt(sl, "solution summary — from design to working prototype",
        0.55, 2.3, 6.3, 0.46, size=13, color=WHITE)

    txt(sl, "PRJ-NS-7421  |  Northstar Civic Group  |  MSA-NS-2024-0418  |  Cycle 2026-04",
        0.45, 2.9, 9, 0.35, size=10, color=DARK_GRAY)

    sections = ["Problem statement", "Architecture", "Classification",
                "Design decisions", "Edge cases", "KPI", "Tech stack"]
    for i, s in enumerate(sections):
        x = 0.45 + i * 1.36
        rect(sl, x, 3.45, 1.28, 0.32, fill=LIGHT_PURPLE)
        txt(sl, s, x + 0.05, 3.47, 1.18, 0.28, size=8, color=DEEP_PURPLE, align=PP_ALIGN.CENTER)

    txt(sl, "Pod 1  |  Meridian Atlas Partners  |  Hackathon Day 2",
        0.45, 5.0, 7, 0.35, size=10, color=WHITE)


# ── Slide 2 — Problem statement ───────────────────────────────────────────────

def slide_problem(prs):
    sl = blank(prs)
    chrome(sl, "Problem statement")

    rect(sl, 0.35, 0.85, 9.3, 0.5, fill=LIGHT_PURPLE)
    txt(sl,
        "An agentic pipeline that reads every backup document line by line, applies contract rules "
        "and institutional memory to classify and auto-resolve expense exceptions, and presents a "
        "pre-populated decision queue — so the analyst reviews only what the machine cannot decide.",
        0.45, 0.87, 9.1, 0.48, size=10, color=DEEP_PURPLE)

    # What the analyst does today
    rect(sl, 0.35, 1.5, 4.35, 0.32, fill=DEEP_PURPLE)
    txt(sl, "What the analyst does today — manually", 0.45, 1.52, 4.15, 0.28,
        size=10.5, bold=True, color=WHITE)

    manual_steps = [
        "Verify every expense against its backup document",
        "Apply contract caps, alcohol exclusions, markup, FX conversion",
        "Follow Project Lead instructions that may override the contract",
        "Flag exceptions, seek resolutions, record decisions",
    ]
    for j, s in enumerate(manual_steps):
        rect(sl, 0.35, 1.87 + j * 0.36, 0.06, 0.28, fill=PURPLE)
        txt(sl, s, 0.48, 1.87 + j * 0.36, 4.2, 0.3, size=9.5, color=BLACK)

    # Scope decision
    rect(sl, 4.95, 1.5, 4.7, 0.32, fill=PURPLE)
    txt(sl, "Scope decision", 5.05, 1.52, 4.5, 0.28, size=10.5, bold=True, color=WHITE)

    rect(sl, 4.95, 1.87, 4.7, 0.9, fill=LIGHT_GRAY)
    txt(sl, "Expenses only — not labour",
        5.05, 1.9, 4.5, 0.28, size=11, bold=True, color=BLACK)
    txt(sl, "Labour reconciliation is a rate-table lookup.\nExpenses require document intelligence — "
        "they are where verification is most complex and error risk is highest.",
        5.05, 2.22, 4.5, 0.52, size=9.5, color=DARK_GRAY)

    # Three root causes
    causes = [
        ("Manual volume", "~27 expense transactions per cycle. Every receipt opened, every amount checked by hand."),
        ("Fragmented memory", "Prior decisions in email, Teams, individual recall. No structure. Instructions contradict across cycles."),
        ("False CLEAN risk", "Document totals match SAP, yet contain non-reimbursable items (alcohol, personal charges) embedded inside."),
    ]
    for i, (title, body) in enumerate(causes):
        x = 0.35 + i * 3.18
        rect(sl, x, 3.5, 3.05, 0.32, fill=DARK_PURPLE)
        txt(sl, title, x + 0.08, 3.52, 2.88, 0.28, size=10, bold=True, color=WHITE)
        txt(sl, body,  x + 0.08, 3.87, 2.88, 0.65, size=9,  color=BLACK)


# ── Slide 3 — Four-agent pipeline ─────────────────────────────────────────────

def slide_architecture(prs):
    sl = blank(prs)
    chrome(sl, "Solution architecture — four-agent pipeline")

    # Pipeline boxes
    agents = [
        ("Ingest\nagent", "Loads SAP CSV,\ncontract, documents,\nprior exceptions"),
        ("Extract + Match\nValidate + Classify\n(Claude opus-4-8)", "Line-level doc extraction\nContract rule check\nClassification output"),
        ("Exception\ntriage agent\n(async)", "Matches flags to\nprior recurring\nresolutions"),
        ("Analyst\nreview\n(Streamlit)", "Exception queue\nPre-populated decisions\nAudit trail capture"),
    ]
    fills = [DEEP_PURPLE, PURPLE, DARK_PURPLE, DEEP_PURPLE]
    for i, ((name, detail), fill) in enumerate(zip(agents, fills)):
        x = 0.3 + i * 2.4
        rect(sl, x, 0.88, 2.15, 0.98, fill=fill)
        txt(sl, name, x + 0.08, 0.9, 1.99, 0.6, size=9.5, bold=True, color=WHITE)
        if i < 3:
            rect(sl, x + 2.17, 1.2, 0.2, 0.3, fill=MID_GRAY)
            txt(sl, ">", x + 2.17, 1.17, 0.2, 0.3, size=13, bold=True,
                color=DARK_GRAY, align=PP_ALIGN.CENTER)
        rect(sl, x, 1.9, 2.15, 0.75, fill=LIGHT_GRAY)
        txt(sl, detail, x + 0.08, 1.93, 1.99, 0.68, size=8.5, color=BLACK)

    # Match key
    rect(sl, 0.3, 2.8, 4.5, 0.28, fill=MID_GRAY)
    txt(sl, "Match key — how transactions link to documents", 0.4, 2.82, 4.3, 0.24,
        size=9.5, bold=True, color=BLACK)
    rect(sl, 0.3, 3.12, 4.5, 0.72, fill=LIGHT_GRAY)
    txt(sl, 'SAP note field: "Receipt: RC-001" / "Vendor invoice: VI-002" / "Mileage log: ML-001"\n'
           'Pipeline parses first RC-* / VI-* / ML-* reference as the primary document link.\n'
           'Amount > $25 with no parseable reference is automatically flagged.',
        0.4, 3.15, 4.3, 0.66, size=8.5, color=BLACK)

    # Architecture decisions
    rect(sl, 5.05, 2.8, 4.6, 0.28, fill=DARK_PURPLE)
    txt(sl, "Key architecture decisions", 5.15, 2.82, 4.4, 0.24,
        size=9.5, bold=True, color=WHITE)

    decisions = [
        ("Separate Ingest agent", "Add new data sources without touching downstream agents"),
        ("Combined Classify agent", "Single Claude call with full context — less latency, better analysis"),
        ("Triage is async", "Classification not blocked; clean lines release immediately"),
        ("claude-opus-4-8", "Adaptive thinking — best-in-class for line-level doc extraction"),
    ]
    for j, (d, r) in enumerate(decisions):
        y = 3.12 + j * 0.38
        bg = LIGHT_PURPLE if j % 2 == 0 else WHITE
        rect(sl, 5.05, y, 4.6, 0.36, fill=bg)
        txt(sl, d, 5.15, y + 0.02, 1.5, 0.32, size=8.5, bold=True, color=DEEP_PURPLE)
        txt(sl, r, 6.7,  y + 0.02, 2.88, 0.32, size=8.5, color=BLACK)


# ── Slide 4 — Classification scheme ──────────────────────────────────────────

def slide_classification(prs):
    sl = blank(prs)
    chrome(sl, "Classification scheme")

    classes = [
        ("CLEAN",       DEEP_PURPLE,  WHITE,
         "Amount matches receipt, no policy violation.",
         "RC-001 flight outbound — SAP $412.80 matches receipt."),
        ("FLAG",        PURPLE,       WHITE,
         "Policy violation, amount mismatch, or ambiguity needing analyst decision.",
         "RC-007 team dinner — alcohol embedded in receipt total."),
        ("EXEMPT",      DARK_PURPLE,  WHITE,
         "No receipt required by contract — per diem or under-$25 threshold.",
         "TX-0030 site parking $18.00 — under $25 threshold."),
        ("MISSING_DOC", MID_GRAY,     BLACK,
         "Document referenced in SAP note but not found in the document repository.",
         "RC-004, RC-005, RC-006, RC-008 to RC-011 — absent from repo."),
        ("ORPHAN",      LIGHT_GRAY,   DARK_GRAY,
         "Document exists in repository with no matching SAP transaction.",
         "RC-019 Harbor Fuel Stop receipt — no SAP transaction found."),
        ("UNREADABLE",  BLACK,        WHITE,
         "Document cannot be parsed — corrupted scan, OCR failure, or malformed agent response.",
         "RC-018 — image too noisy for OCR; engine returned confidence 0.21."),
    ]

    for i, (status, bg, fg, meaning, example) in enumerate(classes):
        row = i // 3
        col = i % 3
        x = 0.3 + col * 3.22
        y = 0.88 + row * 2.18

        rect(sl, x, y, 3.08, 0.42, fill=bg)
        txt(sl, status, x + 0.1, y + 0.05, 2.88, 0.35, size=13, bold=True,
            color=fg, align=PP_ALIGN.LEFT)

        rect(sl, x, y + 0.44, 3.08, 0.75, fill=LIGHT_GRAY)
        txt(sl, meaning, x + 0.1, y + 0.47, 2.88, 0.38, size=9, bold=False, color=BLACK)

        rect(sl, x, y + 1.21, 3.08, 0.72, fill=WHITE, line=MID_GRAY, lw=0.5)
        txt(sl, "Example:", x + 0.1, y + 1.24, 2.88, 0.2, size=8, bold=True, color=DARK_GRAY)
        txt(sl, example,   x + 0.1, y + 1.44, 2.88, 0.44, size=8.5, italic=True, color=BLACK)


# ── Slide 5 — Key design decisions ───────────────────────────────────────────

def slide_decisions(prs):
    sl = blank(prs)
    chrome(sl, "Key design decisions")

    decisions = [
        {
            "num": "1",
            "title": "Line-level extraction",
            "tag": "Most critical",
            "body": ("Documents are analysed line by line — not as a total. "
                     "RC-003 hotel folio: dinner $72.40 contains $3.90 wine. "
                     "RC-007 team dinner: total $126.85 contains House red x2 ($18). "
                     "Both match SAP totals. Both flag as non-reimbursable only at line level."),
            "verdict": "Prevents false CLEAN — the biggest single failure mode",
        },
        {
            "num": "2",
            "title": "Contract as first-class input",
            "tag": "Every call",
            "body": ("MSA-NS-2024-0418 is passed to Claude on every classification call. "
                     "Detects missing 8% subcontractor markup (VI-002: $2,400 cost → $2,592 billable). "
                     "Applies meal cap $90/day, lodging cap $275/$195, travel rate 50%, alcohol exclusion."),
            "verdict": "No rule is applied from memory — always from the contract",
        },
        {
            "num": "3",
            "title": "Exception store replaces channel memory",
            "tag": "Institutional memory",
            "body": ("Prior resolutions with instruction_recurring = Y are standing rules. "
                     "7 of 10 resolutions are recurring — providing the ≥55% auto-resolvable KPI. "
                     "Replaces email / Teams as the source of institutional knowledge. "
                     "Non-recurring resolutions are informational only."),
            "verdict": "Structured, auditable, project-scoped — not in an inbox",
        },
        {
            "num": "4",
            "title": "Cross-project contamination guard",
            "tag": "Data integrity",
            "body": ("EX-2026-0327 is a recurring resolution from PRJ-OTHER-9912 with different rules. "
                     "Despite instruction_recurring = Y it must never apply to PRJ-NS-7421. "
                     "Pipeline hard-filters prior resolutions by project_id before any Claude call."),
            "verdict": "Wrong-project rules silently cause billing errors",
        },
        {
            "num": "5",
            "title": "Immutable audit trail",
            "tag": "2-year retention",
            "body": ("Every pipeline run and analyst decision appended to audit_trail.json. "
                     "Nothing deleted or overwritten. Captures: transaction ID, action, adjusted amount, "
                     "analyst name, timestamp, override flag, full reasoning. "
                     "Analyst must not need to re-open SAP — exception record is self-contained."),
            "verdict": "Full traceability from SAP transaction to invoice line",
        },
        {
            "num": "6",
            "title": "SAP is the system of record",
            "tag": "Integration boundary",
            "body": ("The agent integrates with SAP — it does not replace it. "
                     "SAP amounts are the reference point; agent identifies discrepancies and suggests adjustments. "
                     "No SAP write-backs without analyst approval. "
                     "Project Lead workflow remains undisturbed."),
            "verdict": "Risk-free: no unauthorised financial changes",
        },
    ]

    for i, d in enumerate(decisions):
        row = i // 3
        col = i % 3
        x = 0.3 + col * 3.22
        y = 0.88 + row * 2.2

        # Number badge
        rect(sl, x, y, 0.36, 0.36, fill=PURPLE)
        txt(sl, d["num"], x, y, 0.36, 0.36, size=14, bold=True,
            color=WHITE, align=PP_ALIGN.CENTER)

        # Title + tag
        rect(sl, x + 0.38, y, 2.68, 0.36, fill=LIGHT_GRAY)
        txt(sl, d["title"], x + 0.46, y + 0.02, 1.8, 0.32, size=9.5, bold=True, color=BLACK)
        rect(sl, x + 2.3, y + 0.04, 0.72, 0.28, fill=DEEP_PURPLE)
        txt(sl, d["tag"], x + 2.3, y + 0.04, 0.72, 0.28, size=7,
            bold=True, color=WHITE, align=PP_ALIGN.CENTER)

        # Body
        rect(sl, x, y + 0.38, 3.06, 1.3, fill=WHITE, line=LIGHT_GRAY, lw=0.5)
        txt(sl, d["body"], x + 0.08, y + 0.41, 2.9, 1.0, size=8.5, color=BLACK)

        # Verdict
        rect(sl, x, y + 1.7, 3.06, 0.32, fill=LIGHT_PURPLE)
        txt(sl, d["verdict"], x + 0.08, y + 1.72, 2.9, 0.28, size=8,
            bold=True, color=DEEP_PURPLE, italic=True)


# ── Slide 6 — Edge cases ─────────────────────────────────────────────────────

def slide_edge_cases(prs):
    sl = blank(prs)
    chrome(sl, "Edge cases surfaced in the dataset")

    cases = [
        ("TX-0025 / RC-003", "Hotel folio: dinner $72.40 contains $3.90 wine", "FLAG", "Alcohol line extracted — non-reimbursable"),
        ("TX-0029 / RC-007", "Team dinner: SAP $46.20 vs receipt $126.85; alcohol included", "FLAG", "Prior EX-2025-0911 auto-suggests REJECT"),
        ("TX-0037",          "Drafting supplies $67.30 — no receipt on file", "FLAG", "Auto REJECT via EX-2025-1003 (30-day rule)"),
        ("TX-0040 / RC-012", "Hotel $310 vs $195 elsewhere cap", "FLAG", "Prior EX-2025-0828 auto-approves (PL confirmed)"),
        ("TX-0043 / RC-013", "Client dinner $118 vs $90/day meal cap", "FLAG", "PL note on file — escalate for confirmation"),
        ("TX-0044 / VI-002", "Subcontractor $2,400 — 8% markup not applied", "FLAG", "Auto ADJUST to $2,592 via EX-2025-1129"),
        ("TX-0048 / RC-015", "Receipt issued in CAD on a USD-billed project", "FLAG", "Auto ADJUST with FX note via EX-2026-0314"),
        ("TX-0050 / RC-017", "Personal laundry $22 — non-reimbursable category", "FLAG", "Personal item — contract excludes"),
        ("RC-019",           "Harbor Fuel Stop receipt in repo, no SAP transaction", "ORPHAN", "Escalate: missing transaction or irrelevant doc"),
        ("RC-004 to RC-011", "7 receipts referenced in SAP, absent from document repo", "MISSING_DOC", "Retrieve before billing can proceed"),
        ("RC-018",           "Corrupted scan — OCR confidence 0.21, no readable content", "UNREADABLE", "Flag for analyst — do not infer amounts from noise"),
        ("E-7702",           "Expenses this cycle but no timecard entries in extract", "SYSTEM FLAG", "Surfaced via draft invoice — analyst must verify"),
    ]

    status_colors = {
        "FLAG": PURPLE, "ORPHAN": DARK_PURPLE, "MISSING_DOC": DARK_GRAY,
        "UNREADABLE": BLACK, "SYSTEM FLAG": MID_GRAY,
    }

    # Column headers
    headers = ["Transaction / Doc", "Issue", "Status", "Agent behaviour"]
    col_w = [1.9, 2.8, 1.1, 3.85]
    col_x = [0.2, 2.12, 4.94, 6.06]

    for ci, (h, cx, cw) in enumerate(zip(headers, col_x, col_w)):
        rect(sl, cx, 0.88, cw, 0.3, fill=DEEP_PURPLE)
        txt(sl, h, cx + 0.06, 0.9, cw - 0.1, 0.26, size=8.5, bold=True,
            color=WHITE, align=PP_ALIGN.LEFT)

    for ri, (tx_id, issue, status, behaviour) in enumerate(cases):
        y = 1.2 + ri * 0.355
        bg = LIGHT_GRAY if ri % 2 == 0 else WHITE
        row_vals = [tx_id, issue, status, behaviour]
        for ci, (val, cx, cw) in enumerate(zip(row_vals, col_x, col_w)):
            if ci == 2:
                sc = status_colors.get(status, MID_GRAY)
                rect(sl, cx, y, cw, 0.33, fill=sc)
                txt(sl, val, cx + 0.04, y + 0.04, cw - 0.06, 0.26,
                    size=7.5, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
            else:
                rect(sl, cx, y, cw, 0.33, fill=bg)
                txt(sl, val, cx + 0.06, y + 0.05, cw - 0.1, 0.26,
                    size=8, color=BLACK)


# ── Slide 7 — Auto-resolvable KPI ────────────────────────────────────────────

def slide_kpi(prs):
    sl = blank(prs)
    chrome(sl, "Auto-resolvable KPI — the institutional memory")

    # KPI metric
    rect(sl, 0.3, 0.88, 2.5, 1.1, fill=PURPLE)
    txt(sl, "55%+", 0.3, 0.9, 2.5, 0.65, size=40, bold=True,
        color=WHITE, align=PP_ALIGN.CENTER)
    txt(sl, "of exceptions auto-resolved\nwithout analyst input",
        0.3, 1.55, 2.5, 0.4, size=9, color=WHITE, align=PP_ALIGN.CENTER)

    txt(sl, "Target met from day one — 6 of 7 recurring prior resolutions map "
           "directly to exceptions in cycle 2026-04. No new data collection required.",
        2.95, 0.88, 6.9, 0.55, size=10, color=BLACK)

    rect(sl, 2.95, 1.48, 6.9, 0.04, fill=LIGHT_PURPLE)

    # 7 exception cards
    exceptions = [
        ("EX-2025-0828", "EXPENSE_OVERLIMIT", "Hotel over-cap at coastal site", "APPROVE — no alternative available (PL confirmed)"),
        ("EX-2025-0911", "EXPENSE_DISALLOWED", "Alcohol on any receipt", "REJECT — not reimbursable under any circumstance"),
        ("EX-2025-1003", "MISSING_RECEIPT",    "Expense without receipt on file", "REJECT — hold 30 days, then not billable"),
        ("EX-2025-1129", "SUBCONTRACTOR_MARKUP","Subcontractor invoice, markup absent", "ADJUST — apply 8% per MSA section 3"),
        ("EX-2026-0203", "MISCODED_TIME",       "Internal time coded to client project", "REMOVE — not billable to client"),
        ("EX-2026-0314", "FOREIGN_CURRENCY",    "Receipt in foreign currency (CAD)", "ADJUST — convert at receipt-date FX"),
        ("EX-2025-0712", "RATE_OVERRIDE",       "Off-hours rate claimed (none exists)", "BILL at standard role rate"),
    ]

    for i, (ex_id, ex_type, desc, resolution) in enumerate(exceptions):
        row = i // 4
        col = i % 4
        x = 0.3  + col * 2.42
        y = 1.62 + row * 1.72
        if i == 6:
            x = 0.3 + 1.5 * 2.42   # center last card in second row

        rect(sl, x, y, 2.3, 0.28, fill=DEEP_PURPLE)
        txt(sl, ex_id, x + 0.06, y + 0.03, 1.5, 0.22, size=8, bold=True, color=WHITE)
        rect(sl, x + 1.58, y + 0.03, 0.68, 0.22, fill=PURPLE)
        txt(sl, "recurring", x + 1.58, y + 0.03, 0.68, 0.22,
            size=6.5, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

        rect(sl, x, y + 0.3, 2.3, 0.26, fill=LIGHT_GRAY)
        txt(sl, ex_type, x + 0.06, y + 0.32, 2.18, 0.22, size=7.5, bold=True, color=DARK_GRAY)

        rect(sl, x, y + 0.58, 2.3, 0.5, fill=WHITE, line=MID_GRAY, lw=0.5)
        txt(sl, desc, x + 0.06, y + 0.6, 2.18, 0.22, size=8, color=BLACK)
        rect(sl, x, y + 0.82, 2.3, 0.26, fill=LIGHT_PURPLE)
        txt(sl, resolution, x + 0.06, y + 0.84, 2.18, 0.22, size=7.5, bold=True, color=DEEP_PURPLE)


# ── Slide 8 — Technical stack + what was built ────────────────────────────────

def slide_stack(prs):
    sl = blank(prs)
    chrome(sl, "What was built — technical stack")

    # Stack table left
    rect(sl, 0.3, 0.88, 4.5, 0.3, fill=DEEP_PURPLE)
    txt(sl, "Technical stack", 0.4, 0.9, 4.3, 0.26, size=10, bold=True, color=WHITE)

    stack = [
        ("Language",  "Python 3.9"),
        ("UI",        "Streamlit 1.50"),
        ("AI model",  "claude-opus-4-8 · adaptive thinking · streaming"),
        ("SDK",       "Anthropic Python SDK"),
        ("State",     "JSON files — results, decisions, audit trail"),
        ("Data",      "Synthetic dataset — shared/appendix-sample-data/"),
        ("Tests",     "37 tests · unittest + mock · custom runner"),
    ]
    for j, (component, choice) in enumerate(stack):
        y = 1.22 + j * 0.38
        bg = LIGHT_GRAY if j % 2 == 0 else WHITE
        rect(sl, 0.3, y, 4.5, 0.36, fill=bg)
        txt(sl, component, 0.4,  y + 0.06, 1.1, 0.26, size=9, bold=True, color=DEEP_PURPLE)
        txt(sl, choice,    1.55, y + 0.06, 3.2, 0.26, size=9, color=BLACK)

    # Files right
    rect(sl, 5.1, 0.88, 4.55, 0.3, fill=PURPLE)
    txt(sl, "Files delivered", 5.2, 0.9, 4.35, 0.26, size=10, bold=True, color=WHITE)

    files = [
        ("app.py",             "Streamlit UI — Dashboard, Exception Queue, All Transactions, Audit Trail"),
        ("pipeline.py",        "3-agent pipeline: Ingest, Classify (Claude), Triage (async)"),
        ("data_loader.py",     "SAP CSV, contract, 15 backup documents, 10 prior exceptions"),
        ("test_pipeline.py",   "37 automated tests — data loading, rule-based, mocked Claude, state"),
        ("requirements.txt",   "anthropic · streamlit · pandas"),
        ("SOLUTION-SUMMARY.md","Full design rationale and decision log"),
        ("RDE-PLAYBOOK-OUTPUT.md", "Structured output mapped to 6 RDE playbook phases"),
    ]
    for j, (fname, desc) in enumerate(files):
        y = 1.22 + j * 0.38
        bg = LIGHT_PURPLE if j % 2 == 0 else WHITE
        rect(sl, 5.1, y, 4.55, 0.36, fill=bg)
        txt(sl, fname, 5.2,  y + 0.06, 1.55, 0.26, size=8.5, bold=True, color=DEEP_PURPLE)
        txt(sl, desc,  6.8,  y + 0.06, 2.78, 0.26, size=8.5, color=BLACK)

    # Run command footer
    rect(sl, 0.3, 4.98, 9.35, 0.42, fill=DEEP_PURPLE)
    txt(sl, "Run:  cd pod1  |  export ANTHROPIC_API_KEY=...  |  streamlit run app.py  →  http://localhost:8501",
        0.45, 5.02, 9.1, 0.34, size=9.5, bold=True, color=WHITE)


# ── Slide 9 — Constraints + summary ──────────────────────────────────────────

def slide_summary(prs):
    sl = blank(prs)
    chrome(sl, "Summary — what the prototype proves")

    proofs = [
        ("Agent classifies accurately",
         "Claude correctly identifies alcohol embedded in RC-003 and RC-007. "
         "Flags missing 8% markup in VI-002. Surfaces RC-019 orphan automatically."),
        ("55%+ auto-resolvable KPI met",
         "6 of 7 recurring resolutions map directly to exceptions in this cycle. "
         "No new data collection required before the first run."),
        ("False CLEAN rate: zero",
         "Line-level extraction catches non-reimbursable items that total-level "
         "checks miss. The biggest single failure mode is eliminated by design."),
        ("Human stays in control",
         "Every classification is a suggestion. Analyst confirms or overrides. "
         "Every decision is logged immutably. No SAP write-back without approval."),
        ("Architecture generalises",
         "Swap the Ingest agent's data source to add a new project. "
         "Classify and Triage agents are untouched. Labour and multi-project "
         "batch are additive — no redesign required."),
    ]

    for i, (title, body) in enumerate(proofs):
        row = i // 3
        col = i % 3
        if i < 3:
            x, y, w = 0.3 + col * 3.22, 0.88, 3.08
        else:
            x, y, w = 0.3 + (i - 3) * 4.85, 2.72, 4.62

        rect(sl, x, y, w, 0.3, fill=PURPLE)
        txt(sl, title, x + 0.08, y + 0.03, w - 0.14, 0.24, size=9.5, bold=True, color=WHITE)
        rect(sl, x, y + 0.32, w, 0.88, fill=LIGHT_GRAY)
        txt(sl, body, x + 0.08, y + 0.35, w - 0.14, 0.82, size=9, color=BLACK)

    # Constraints banner
    rect(sl, 0.3, 4.28, 9.35, 0.3, fill=DARK_GRAY)
    txt(sl, "Constraints respected:", 0.4, 4.3, 1.6, 0.26, size=8.5, bold=True, color=WHITE)
    constraints = [
        "Synthetic data only", "SAP = system of record",
        "No SAP write-backs", "Project Lead workflow undisturbed",
    ]
    for i, c in enumerate(constraints):
        x = 2.1 + i * 1.88
        txt(sl, f"  {c}", x, 4.3, 1.82, 0.26, size=8.5, color=WHITE)

    # Bug fix note
    rect(sl, 0.3, 4.65, 9.35, 0.66, fill=LIGHT_PURPLE)
    txt(sl, "Bug found and fixed during sprint:",
        0.4, 4.67, 2.8, 0.22, size=8.5, bold=True, color=DEEP_PURPLE)
    txt(sl, "Pipeline Error: 'classification'  —  max_tokens=2048 caused adaptive "
           "thinking to exhaust the token budget; Claude JSON output was truncated; "
           "KeyError on missing key. Fixed: raised to 8192, added explicit guard after JSON parse, "
           "changed r[\"classification\"] to r.get(\"classification\").",
        0.4, 4.89, 9.1, 0.38, size=8.5, color=BLACK)


# ── Build ─────────────────────────────────────────────────────────────────────

def build():
    prs = new_prs()
    slide_cover(prs)
    slide_problem(prs)
    slide_architecture(prs)
    slide_classification(prs)
    slide_decisions(prs)
    slide_edge_cases(prs)
    slide_kpi(prs)
    slide_stack(prs)
    slide_summary(prs)
    prs.save(str(OUT_PATH))
    print(f"Saved: {OUT_PATH}")


if __name__ == "__main__":
    build()
